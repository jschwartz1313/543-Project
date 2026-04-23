from __future__ import annotations

import io
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

ROOT = pathlib.Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / "outputs" / "figures"
OUT = ROOT / "outputs" / "ercot_winter_storm_analysis.pptx"

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1a, 0x3a, 0x5c)
RED    = RGBColor(0xc8, 0x39, 0x2b)
BLUE   = RGBColor(0x2a, 0x64, 0x96)
LGRAY  = RGBColor(0xf2, 0xf6, 0xfc)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
WHITE  = RGBColor(0xff, 0xff, 0xff)
BLACK  = RGBColor(0x11, 0x11, 0x11)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height


# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def rect(slide, l, t, w, h, fill: RGBColor | None = None) -> object:
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def add_text(slide, text: str, l, t, w, h, *,
             size: int = 18, bold: bool = False, color: RGBColor = BLACK,
             align=PP_ALIGN.LEFT, wrap: bool = True) -> None:
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def header_bar(slide, title: str) -> None:
    rect(slide, 0, 0, W, Inches(0.85), fill=NAVY)
    add_text(slide, title,
             Inches(0.35), Inches(0.12), Inches(12.5), Inches(0.65),
             size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def accent_line(slide) -> None:
    bar = rect(slide, 0, Inches(0.85), W, Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()


def add_figure(slide, fname: str, l, t, w, h) -> None:
    path = FIG_DIR / fname
    slide.shapes.add_picture(str(path), l, t, w, h)


def footnote(slide, text: str) -> None:
    add_text(slide, text,
             Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.35),
             size=9, color=DGRAY, align=PP_ALIGN.LEFT)


def bullet_box(slide, items: list[str], l, t, w, h, size: int = 14) -> None:
    txb = slide.shapes.add_textbox(l, t, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"  •  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = BLACK


def simple_table(slide, headers: list[str], rows: list[list[str]],
                 l, t, w, h) -> None:
    cols = len(headers)
    col_w = w // cols
    row_h = h // (len(rows) + 1)

    # header row
    for ci, hdr in enumerate(headers):
        bg = rect(slide, l + ci * col_w, t, col_w, row_h, fill=NAVY)
        add_text(slide, hdr,
                 l + ci * col_w + Inches(0.05), t + Inches(0.02),
                 col_w - Inches(0.1), row_h,
                 size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        bg_color = LGRAY if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            rect(slide, l + ci * col_w, t + (ri + 1) * row_h,
                 col_w, row_h, fill=bg_color)
            add_text(slide, cell,
                     l + ci * col_w + Inches(0.05),
                     t + (ri + 1) * row_h + Inches(0.02),
                     col_w - Inches(0.1), row_h,
                     size=10, color=BLACK, align=PP_ALIGN.CENTER)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs: Presentation) -> None:
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, 0, Inches(5.8), W, Inches(1.7), fill=RED)

    add_text(sl, "ERCOT Winter Storm Outage Risk",
             Inches(0.6), Inches(1.4), Inches(12.0), Inches(1.6),
             size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "Utility exposure, financial impact, and a URI-like scenario analysis",
             Inches(0.6), Inches(3.1), Inches(12.0), Inches(0.7),
             size=20, color=RGBColor(0xcc, 0xdd, 0xee), align=PP_ALIGN.CENTER)

    tags = ["EIA-861 Reliability", "NOAA Storm Events", "NOAA GHCN Daily",
            "6 ERCOT Utilities · 2013–2024"]
    tag_str = "   |   ".join(tags)
    add_text(sl, tag_str,
             Inches(0.6), Inches(6.05), Inches(12.0), Inches(0.55),
             size=13, color=WHITE, align=PP_ALIGN.CENTER)


def slide_research_question(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Research Question")
    accent_line(sl)

    add_text(sl,
             '"How much can winter storm conditions in ERCOT help explain utility outages, '
             'and what do those outages imply for delivery-rate pressure and lost revenue risk?"',
             Inches(0.5), Inches(1.0), Inches(12.3), Inches(1.1),
             size=15, color=NAVY, align=PP_ALIGN.CENTER)

    rect(sl, Inches(0.35), Inches(2.2), Inches(5.9), Inches(4.4), fill=LGRAY)
    add_text(sl, "Sub-questions",
             Inches(0.5), Inches(2.3), Inches(5.6), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "Which territories have the largest winter storm exposure?",
        "How do outage metrics vary across utilities over time?",
        "How much delivery revenue is at risk when outage duration rises?",
        "What would a URI-like profile imply for 2024 conditions?",
    ], Inches(0.5), Inches(2.75), Inches(5.6), Inches(3.6))

    rect(sl, Inches(6.6), Inches(2.2), Inches(6.38), Inches(4.4), fill=LGRAY)
    add_text(sl, "Scope",
             Inches(6.75), Inches(2.3), Inches(6.1), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "Geography: ERCOT T&D utilities in Texas",
        "Storm exposure: 2000–2025",
        "Outage metrics: 2013–2024 (EIA-861)",
        "Revenue metrics: 2020–2024 (EIA-861)",
    ], Inches(6.75), Inches(2.75), Inches(6.1), Inches(3.6))


def slide_utilities(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "The Six ERCOT Delivery Utilities")
    accent_line(sl)

    headers = ["Utility", "Territory"]
    rows = [
        ["AEP Texas Central",    "South Texas (Corpus Christi, Rio Grande Valley)"],
        ["AEP Texas North",      "West & Panhandle Texas (Abilene, Lubbock area)"],
        ["CenterPoint Energy",   "Greater Houston metro"],
        ["Nueces Electric Coop", "Coastal Bend (Corpus Christi fringe)"],
        ["Oncor Electric",       "North & Central Texas (DFW, West Texas)"],
        ["Texas-New Mexico Power", "Scattered suburban corridors statewide"],
    ]
    simple_table(sl, headers, rows,
                 Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.4))
    footnote(sl, "Together these utilities serve the majority of the ERCOT transmission and distribution footprint.")


def slide_data_sources(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Data Sources")
    accent_line(sl)

    boxes = [
        ("EIA Form 861",
         ["Annual reliability schedules (SAIDI, SAIFI)",
          "Service territory county mapping",
          "Delivery revenue & MWh sales"]),
        ("NOAA Storm Events",
         ["County-level winter-season storm records",
          "Property & crop damage estimates",
          "Direct deaths & injuries"]),
        ("NOAA GHCN Daily",
         ["One representative station per utility",
          "Freeze days, sub-20°F days",
          "Snowfall, precip, min temperature"]),
    ]
    box_w = Inches(4.0)
    for i, (title, items) in enumerate(boxes):
        lx = Inches(0.35) + i * Inches(4.33)
        rect(sl, lx, Inches(1.1), box_w, Inches(5.5), fill=LGRAY)
        add_text(sl, title, lx + Inches(0.15), Inches(1.2),
                 box_w - Inches(0.3), Inches(0.45),
                 size=15, bold=True, color=NAVY)
        bullet_box(sl, items,
                   lx + Inches(0.1), Inches(1.75),
                   box_w - Inches(0.2), Inches(4.5),
                   size=13)
    footnote(sl, "All raw extracts saved to data/raw/; merged outputs to data/processed/")


def slide_method(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Method")
    accent_line(sl)

    steps = [
        ("1", "Build utility geography",
         "Map ERCOT delivery utilities to Texas counties via EIA service territory data."),
        ("2", "Build storm & weather exposure",
         "Aggregate NOAA winter storm events and daily weather severity to the utility-year level."),
        ("3", "Join utility performance data",
         "Merge annual outage metrics (EIA reliability) and delivery revenue (EIA delivery companies)."),
        ("4", "Fit exploratory outage model",
         "Linear regression of SAIDI on winter severity and storm exposure features."),
        ("5", "Simulate URI-like scenario",
         "Apply 2021 winter feature values to 2024 utility baselines to estimate incremental lost revenue."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        ty = Inches(1.1) + i * Inches(1.2)
        rect(sl, Inches(0.35), ty, Inches(0.6), Inches(0.6), fill=RED)
        add_text(sl, num,
                 Inches(0.35), ty + Inches(0.05), Inches(0.6), Inches(0.55),
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(sl, title,
                 Inches(1.1), ty, Inches(11.5), Inches(0.38),
                 size=14, bold=True, color=NAVY)
        add_text(sl, desc,
                 Inches(1.1), ty + Inches(0.38), Inches(11.5), Inches(0.75),
                 size=12, color=BLACK)


def slide_figure(prs: Presentation, title: str, fname: str, note: str) -> None:
    sl = blank_slide(prs)
    header_bar(sl, title)
    accent_line(sl)
    add_figure(sl, fname,
               Inches(0.5), Inches(0.95), Inches(12.33), Inches(5.85))
    footnote(sl, note)


def slide_uri_callout(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Winter Storm Uri — 2021")
    accent_line(sl)

    headers = ["Utility", "2021 SAIDI (min)", "Normal baseline"]
    rows = [
        ["CenterPoint Energy",      "2,366", "~150"],
        ["AEP Texas Central",       "1,901", "~170"],
        ["Nueces Electric Coop",    "1,524", "~170"],
        ["AEP Texas North",         "1,137", "~110"],
        ["Oncor",                   "  559", "~100"],
        ["Texas-New Mexico Power",  "  204", "~105"],
    ]
    simple_table(sl, headers, rows,
                 Inches(1.0), Inches(1.1), Inches(11.3), Inches(5.0))
    footnote(sl,
             "Uri caused statewide generation failures in February 2021, producing sustained outages across nearly all utilities. "
             "Normal baseline is the approximate SAIDI-without-MED for each utility.")


def slide_model_summary(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Exploratory Outage Model")
    accent_line(sl)

    rect(sl, Inches(0.35), Inches(1.1), Inches(5.8), Inches(5.8), fill=LGRAY)
    add_text(sl, "Setup", Inches(0.5), Inches(1.2), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "OLS linear regression",
        "Dependent variable: SAIDI with MED",
        "48 utility-year observations",
        "8 winter severity features",
    ], Inches(0.5), Inches(1.65), Inches(5.5), Inches(2.0))

    add_text(sl, "Performance", Inches(0.5), Inches(3.8), Inches(5.5), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "In-sample R² = 0.184",
        "In-sample MAE = 442 minutes",
    ], Inches(0.5), Inches(4.25), Inches(5.5), Inches(1.5))

    rect(sl, Inches(6.5), Inches(1.1), Inches(6.48), Inches(5.8), fill=LGRAY)
    add_text(sl, "Features", Inches(6.65), Inches(1.2), Inches(6.2), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "Winter event count",
        "Storm property damage (USD)",
        "Freeze days (≤32°F)",
        "Sub-20°F days",
        "Precipitation days",
        "Snow days",
        "Min temperature (°C)",
        "Total precipitation (mm)",
    ], Inches(6.65), Inches(1.65), Inches(6.2), Inches(5.0))

    footnote(sl, "Model is exploratory — the small panel and annual SAIDI metric limit causal claims.")


def slide_scenario_table(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "URI-like 2024 Scenario")
    accent_line(sl)

    add_text(sl,
             "Apply 2021 winter severity feature values to 2024 utility baselines. Estimate incremental SAIDI and lost revenue.",
             Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.5),
             size=13, color=DGRAY)

    headers = ["Utility", "2024 SAIDI", "URI-like SAIDI", "Incremental Lost Revenue"]
    rows = [
        ["Oncor",                "532",   "685",   "+$1,423,255"],
        ["AEP Texas Central",    "229",   "265",   "+$78,667"],
        ["AEP Texas North",      "108",   "254",   "+$64,943"],
        ["Texas-NM Power",     "1,946", "1,892",   "−$42,413"],
        ["CenterPoint",        "4,316", "4,100",   "−$1,310,112"],
        ["Nueces",               "437",    "—",  "missing 2021 weather"],
    ]
    simple_table(sl, headers, rows,
                 Inches(0.35), Inches(1.5), Inches(12.63), Inches(5.2))
    footnote(sl,
             "CenterPoint and TNMP show negative incremental impact because their 2024 SAIDI (driven by Hurricane Beryl / "
             "summer storms) already exceeds what a URI-level winter would predict.")


def slide_limitations(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Limitations")
    accent_line(sl)

    rect(sl, Inches(0.35), Inches(1.1), Inches(6.0), Inches(5.8), fill=LGRAY)
    add_text(sl, "Data", Inches(0.5), Inches(1.2), Inches(5.7), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "Annual SAIDI captures all outage causes — hurricanes and summer storms inflate "
        "\"winter risk\" metrics for some utilities",
        "One weather station per utility is a coarse spatial proxy",
        "NOAA Storm Events classification does not always align with ERCOT winter hazard definitions",
    ], Inches(0.5), Inches(1.65), Inches(5.7), Inches(5.0), size=13)

    rect(sl, Inches(6.7), Inches(1.1), Inches(6.28), Inches(5.8), fill=LGRAY)
    add_text(sl, "Model", Inches(6.85), Inches(1.2), Inches(6.0), Inches(0.4),
             size=14, bold=True, color=NAVY)
    bullet_box(sl, [
        "Small panel (48 obs across 6 utilities)",
        "Low R² (0.18) — winter features explain limited variance in annual SAIDI",
        "OLS assumes linear, additive effects; no interaction terms",
        "Not causal — confounders like infrastructure age are omitted",
    ], Inches(6.85), Inches(1.65), Inches(6.0), Inches(5.0), size=13)


def slide_next_steps(prs: Presentation) -> None:
    sl = blank_slide(prs)
    header_bar(sl, "Next Steps")
    accent_line(sl)

    bullet_box(sl, [
        "Restrict outage metric to winter months only to remove summer storm contamination",
        "Replace single-station weather proxy with a broader spatial weather surface",
        "Add a richer outage/event source that ties more directly to utility-specific winter incidents",
        "Test alternative model forms (fixed effects, log-SAIDI, nonlinear weather terms)",
        "Tighten the financial model to reflect actual utility recovery and rate mechanisms",
    ], Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.8), size=16)


def slide_closing(prs: Presentation) -> None:
    sl = blank_slide(prs)
    rect(sl, 0, 0, W, H, fill=NAVY)
    rect(sl, 0, Inches(5.8), W, Inches(1.7), fill=RED)

    add_text(sl, "ERCOT Winter Storm Outage Risk",
             Inches(0.6), Inches(1.6), Inches(12.0), Inches(1.4),
             size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, "A reproducible analytical baseline combining\nEIA-861, NOAA Storm Events, and NOAA GHCN Daily",
             Inches(0.6), Inches(3.2), Inches(12.0), Inches(1.0),
             size=18, color=RGBColor(0xcc, 0xdd, 0xee), align=PP_ALIGN.CENTER)

    refs = (
        "scripts/download_data.py  ·  scripts/build_analysis.py  ·  "
        "data/processed/ercot_utility_winter_risk_panel.csv  ·  docs/analysis_summary.md"
    )
    add_text(sl, refs,
             Inches(0.6), Inches(6.1), Inches(12.0), Inches(0.5),
             size=11, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main() -> None:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_research_question(prs)
    slide_utilities(prs)
    slide_data_sources(prs)
    slide_method(prs)
    slide_figure(prs,
        "Geographic Winter Storm Exposure",
        "ercot_county_winter_storm_choropleth.png",
        "County-level cumulative NOAA winter storm event counts, 2000-2025. Darker = more events.")
    slide_figure(prs,
        "Utility Outage Trends (SAIDI with MED)",
        "ercot_utility_saidi_timeseries.png",
        "SAIDI with Major Event Days. Spikes reflect large discrete events — notably Winter Storm Uri (2021) "
        "and Hurricane Beryl / summer storms (2024 CenterPoint & TNMP).")
    slide_uri_callout(prs)
    slide_figure(prs,
        "Delivery Revenue per MWh",
        "ercot_delivery_rate_timeseries.png",
        "Delivery revenue ÷ MWh sales from EIA-861 Delivery Company schedules, 2020-2024.")
    slide_model_summary(prs)
    slide_figure(prs,
        "Model Coefficients",
        "ercot_outage_model_coefficients.png",
        "Larger absolute coefficients indicate stronger model association with SAIDI. Interpret with caution given low R².")
    slide_scenario_table(prs)
    slide_figure(prs,
        "URI-like Incremental Lost Revenue",
        "ercot_uri_like_lost_revenue.png",
        "Incremental lost revenue = URI-like predicted lost revenue minus 2024 baseline lost revenue. "
        "Negative values indicate 2024 actual outages already exceed the URI-like prediction.")
    slide_limitations(prs)
    slide_next_steps(prs)
    slide_closing(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT}  ({OUT.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
